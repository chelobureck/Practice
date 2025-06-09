import requests
from pptx import Presentation
import aspose.slides as slides
import aspose.pydrawing as drawing
import os
from fastapi import FastAPI, Request, BackgroundTasks
from fastapi.responses import JSONResponse, FileResponse
import uuid
import base64
from services.file_cleanup import cleanup_presentation_files
import shutil
import logging

app = FastAPI()

def make_pptx_file(plan, filename: str = "presentation.pptx") -> str:
    prs = Presentation()
    for item in plan:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = item.get("title", "")
        slide.placeholders[1].text = item.get("text", "")
        image_url = item.get("image", "")
        if image_url:
            try:
                img_data = requests.get(image_url, timeout=10).content
                img_path = filename + "_temp.png"
                with open(img_path, "wb") as f:
                    f.write(img_data)
                slide.shapes.add_picture(img_path, left=prs.slide_width//3, top=prs.slide_height//3, width=prs.slide_width//3)
                os.remove(img_path)
            except Exception as e:
                print(f"Ошибка загрузки изображения: {e}")
    prs.save(filename)
    return filename

def pptx_to_png(pptx_path, output_dir):
    pres = slides.Presentation(pptx_path)
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    png_files = []
    for i, slide in enumerate(pres.slides):
        out_path = os.path.join(output_dir, f"slide_{i+1}.png")
        image = slide.get_thumbnail(1, 1)
        image.save(out_path, drawing.imaging.ImageFormat.png)
        png_files.append(out_path)
    return png_files

@app.post("/generate_presentation")
async def generate_presentation(request: Request):
    data = await request.json()
    pptx_code = data.get("pptx_code")
    if not pptx_code:
        return JSONResponse({"error": "No pptx code provided"}, status_code=400)

    pres_id = uuid.uuid4().hex
    filename = f"presentation_{pres_id}.pptx"
    pptx_path = make_pptx_file(pptx_code, filename)

    png_dir = f"png_{pres_id}"
    os.makedirs(png_dir, exist_ok=True)
    png_files = pptx_to_png(pptx_path, png_dir)
    images = []
    for png_file in png_files:
        with open(png_file, "rb") as f:
            images.append("data:image/png;base64," + base64.b64encode(f.read()).decode())
    # Удаляем файлы после генерации (асинхронно)
    import threading
    threading.Thread(target=cleanup_presentation_files, args=(pres_id,)).start()
    return JSONResponse({"images": images})

@app.post("/gpt")
async def gpt_pptx(request: Request, background_tasks: BackgroundTasks):
    try:
        data = await request.json()
        history = data.get("history", [])
        if not history:
            return JSONResponse({"error": "No chat history provided"}, status_code=400)

        plan = get_gpt_plan(history)
        if not validate_plan(plan):
            logging.error(f"Invalid plan from GPT: {plan}")
            return JSONResponse({"error": "Invalid structure from GPT"}, status_code=400)

        pres_id = str(uuid.uuid4())
        pres_dir = f"presentations/{pres_id}"
        os.makedirs(pres_dir, exist_ok=True)
        pptx_path = os.path.join(pres_dir, "presentation.pptx")
        make_pptx_file(plan, pptx_path)
        save_last_pptx_path(pptx_path)

        png_dir = os.path.join(pres_dir, "slides")
        os.makedirs(png_dir, exist_ok=True)
        png_files = pptx_to_png(pptx_path, png_dir)
        images = []
        for png_file in png_files:
            with open(png_file, "rb") as f:
                images.append("data:image/png;base64," + base64.b64encode(f.read()).decode())
        # PNG удаляем только после ответа
        background_tasks.add_task(shutil.rmtree, pres_dir, ignore_errors=True)
        return JSONResponse({"images": images, "presentation_id": pres_id})
    except Exception as e:
        logging.error("Ошибка при создании презентации", exc_info=True)
        return JSONResponse({"error": "Failed to create presentation"}, status_code=500)

@app.get("/download_pptx/{pres_id}")
def download_pptx(pres_id: str, background_tasks: BackgroundTasks):
    pptx_path = f"presentations/{pres_id}/presentation.pptx"
    if os.path.exists(pptx_path):
        response = FileResponse(
            path=pptx_path,
            filename="presentation.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        # Удаляем всю папку после скачивания
        background_tasks.add_task(shutil.rmtree, f"presentations/{pres_id}", ignore_errors=True)
        return response
    return JSONResponse({"error": "Файл не найден"}, status_code=404)